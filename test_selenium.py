from bs4 import BeautifulSoup
from selenium import webdriver
from selenium.webdriver.chrome.options import Options
import chromedriver_binary
import time
from selenium.webdriver.common.by import By
#------------------------------------------------#
from pptx import Presentation

prs = Presentation()
title_slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "ハローワールド!"
subtitle.text = "こんにちは、あなたの健康を守ります。"

prs.save('test2.pptx')
#------------------------------------------------#
option = Options()
option.add_argument('--headless')

driver = webdriver.Chrome(options=option)
driver.get('https://www.google.co.jp/')
html = driver.page_source
soup = BeautifulSoup(html, 'html.parser')

# 検索フィールドの取得
query = driver.find_element(by=By.NAME, value='q')

# 検索文字列を入力
query.send_keys('pythonとは')

# 3秒待つ
time.sleep(3)

# 検索ボタンをクリック
button = driver.find_element(by=By.NAME, value='btnK')
button.click()

# 3秒待つ
time.sleep(3)

html = driver.page_source
soup = BeautifulSoup(html, 'html.parser')

ll = [x for x in soup.text.split(' ') if len(x) > 0]
for elem in ll:
    print(prs.slides[0])